import time
from openpyxl import Workbook
from app.models.basemodels import Register_Response_
from app.models.tables import Person as Flask_Person, Spot as Flask_Spot, Spot_Commercial_Info as Flask_Spot_Comm, Spot_Private_Info as Flask_Spot_Int, User, Worksheet_Content as Flask_Worksheet_Content
from sqlalchemy import Column, Date, Float, ForeignKey, Integer, String, create_engine
from app.providers.db_services import get_fornecedores_list, get_fornecedores_list_off_context, get_pontos_by_id_list
from app.providers.s3_services import upload_file_to_s3
from sqlalchemy.orm import declarative_base, Session
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from sqlalchemy.dialects.mysql import JSON
from reportlab.lib.pagesizes import mm
from PIL import UnidentifiedImageError, ImageFile
from reportlab.pdfgen import canvas
from pptx.dml.color import RGBColor
from secrets import token_hex, token_urlsafe
from flask import flash, redirect
from reportlab.lib import colors
from pptx import Presentation
from datetime import date
from pptx.util import Mm
from json import dumps
from worker import conn
from io import BytesIO
from PIL import Image
from rq import Queue
import pandas as pd
from app import db, q
import requests
import os
from sqlalchemy.orm import relationship
import googlemaps

gmaps = googlemaps.Client(key='AIzaSyAdLDJ7M8krlagfM-HOHUxxfgr9N6_nto8')
Base = declarative_base()

class Worksheet_Content(Base):
    __tablename__ = 'worksheet_content'
    id = Column(Integer, autoincrement=True, primary_key=True, nullable=False)
    user_id = Column(Integer)
    title = Column(String(255), nullable=False)
    company = Column(String(255))
    person = Column(String(255))
    content = Column(JSON, nullable=False)
    creation_date = Column(Date, nullable=False)
    image_id = Column(String(255), nullable=False, unique=True)

    def __init__(self, user_id, title, company, person, content, creation_date, image_id):
        self.title = title
        self.user_id = user_id
        self.company = company
        self.person = person
        self.content = content
        self.creation_date = creation_date
        self.image_id = image_id

class Spot(Base):
    __tablename__ = 'spot'
    id = Column(Integer, autoincrement=True, primary_key=True, nullable=False)
    code = Column(String(255), nullable=False)
    address = Column(String(255), nullable=False)
    latitude = Column(Float, nullable=False)
    longitude = Column(Float, nullable=False)
    image_link = Column(String(255), nullable=False)
    include_date = Column(Date, nullable=False)
    reference = Column(String(255))
    district = Column(String(255))
    city = Column(String(255))
    zone = Column(String(255))
    state = Column(String(255))
    country = Column(String(255))
    format = Column(String(255))
    measure = Column(String(255))
    commercial_spot_id = relationship("Spot_Commercial_Info", cascade='all, delete')
    private_spot_id = relationship("Spot_Private_Info", cascade='all, delete')

    def __init__(self, code, address, latitude, longitude, image_link, include_date, reference, district, city, zone, state, country, format, measure):
        self.code = code
        self.address = address
        self.latitude = latitude
        self.longitude = longitude
        self.image_link = image_link
        self.include_date = include_date
        self.reference = reference
        self.district = district
        self.city = city
        self.zone = zone
        self.state = state
        self.country = country
        self.format = format
        self.measure = measure
    
class Spot_Commercial_Info(Base):
    __tablename__ = 'spot_commercial_info'
    id = Column(Integer, autoincrement=True, primary_key=True, nullable=False)
    spot_id = Column(Integer, ForeignKey(Spot.id, ondelete='CASCADE'), nullable=False, unique=True)
    impacto = Column(String(255))
    valor_tabela_comm = Column(String(255))
    valor_negociado_comm = Column(String(255))
    producao = Column(String(255))
    observacoes = Column(String(255))
    outros = Column(JSON(255))
    
    def __init__(self, spot_id, impacto, valor_tabela_comm, valor_negociado_comm, producao, observacoes, outros):
        self.spot_id = spot_id
        self.impacto = impacto
        self.valor_tabela_comm = valor_tabela_comm
        self.valor_negociado_comm = valor_negociado_comm
        self.producao = producao
        self.observacoes = observacoes
        self.outros = outros

class Spot_Private_Info(Base):
    __tablename__ = 'spot_private_info'
    id = Column(Integer, autoincrement=True, primary_key=True, nullable=False)
    spot_id = Column(Integer, ForeignKey(Spot.id, ondelete='CASCADE'), nullable=False, unique=True)
    empresa = Column(String(255))
    valor_negociado_int = Column(String(255))
    custo_liq = Column(String(255))
    medida_int = Column(String(255))
    observacoes = Column(String(255))
    outros = Column(JSON)
    
    def __init__(self, spot_id, empresa, medida_int, valor_negociado_int, custo_liq, observacoes, outros):
        self.spot_id = spot_id
        self.empresa = empresa
        self.medida_int = medida_int
        self.valor_negociado_int = valor_negociado_int
        self.custo_liq = custo_liq
        self.observacoes = observacoes
        self.outros = outros

class Person(Base):
    __tablename__ = 'person'
    id = Column(Integer, autoincrement=True, primary_key=True, nullable=False)
    name = Column(String(255), nullable=False)
    site = Column(String(255))
    person_name = Column(String(255))
    email1 = Column(String(255))
    email2 = Column(String(255))
    tel1 = Column(String(255))
    tel2 = Column(String(255))
    relation_level = Column(Integer)
    person_type = Column(Integer) # 1 = Fornecedor. 2 = Cliente

    def __init__(self, name, site, person_name, email1, email2, tel1, tel2, relation_level, person_type):
        self.name = name
        self.site = site
        self.person_name = person_name
        self.email1 = email1
        self.email2 = email2
        self.tel1 = tel1
        self.tel2 = tel2
        self.relation_level = relation_level
        self.person_type = person_type

engine = create_engine(os.environ['SQLALCHEMY_DATABASE_URI'], echo=False, future=True)
session = Session(engine)

ALLOWED_EXTENSIONS = {'xlsx'}
def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def alternative_id_generator():
    alternative_id = token_hex(16)
    alternative_id_check = User.query.filter(User.alternative_id == alternative_id).first()
    if alternative_id_check != None:
        alternative_id_generator()
    else:
        return alternative_id

def image_id_generator():
    image_id = token_urlsafe(40)
    image_id_db = Flask_Worksheet_Content.query.filter(Flask_Worksheet_Content.image_id == image_id).first()
    if image_id_db:
        image_id_generator()
    return image_id

def flash_dif_languages(lang, es, en, pt):
    if lang == 'es-ar' or lang == 'es':
        return flash(es)
    elif lang == 'en':
        return flash(en)
    else:
        return flash(pt)

def is_in_the_column(coluna, lista):
    for item in lista:
        if coluna.replace(' ', '') in item:
            return True
    return False

def file_validator(arquivo, lang):
    messages = []
    if not arquivo or arquivo.filename == '':
        if lang == 'es' or lang == 'es-ar':
            messages.append('ningun archivo fue cargado.')
        elif lang == 'en':
            messages.append('no files were uploaded.')
        else:
            messages.append('Nenhum arquivo foi enviado.')
        return False, messages
    elif not allowed_file(arquivo.filename):
        if lang == 'es' or lang == 'es-ar':
            messages.append('Formato de archivo no soportado. Debe cargar archivos ".xlsx".')
        elif lang == 'en':
            messages.append('Unsupported file format. You must upload ".xlsx" files.')
        else:
            messages.append('Formato de arquivo não suportado. Você deve enviar arquivos ".xlsx".')
        return False, messages
    return True, messages

def is_float(string):
    try:
        test = float(string)
        return True
    except ValueError as error:
        return False

def pdf_generator(capa, content, image_id, lang, user_id, is_worker):
    try:
        ImageFile.LOAD_TRUNCATED_IMAGES = True
        colunas = content['colunas']
        linhas = content['conteudo']

        endereco_column = None
        latitude_column = None
        longitude_column = None
        codigo_column = None
        foto_column = None
        other_columns = []
        for coluna in colunas:
            if coluna.lower().rsplit(' ')[0] == 'endereco' or coluna.lower().rsplit(' ')[0] == 'endereço' or coluna.lower().rsplit(' ')[0] == 'direccion' or coluna.lower().rsplit(' ')[0] == 'dirección' or coluna.lower().rsplit(' ')[0] == 'ubicacion' or coluna.lower().rsplit(' ')[0] == 'ubicación' or coluna.lower().rsplit(' ')[0] == 'address':
                endereco_column = coluna
            elif coluna.lower().rsplit(' ')[0] == 'latitude' or coluna.lower().rsplit(' ')[0] == 'latitud':
                latitude_column = coluna
            elif coluna.lower().rsplit(' ')[0] == 'longitude' or coluna.lower().rsplit(' ')[0] == 'longitud':
                longitude_column = coluna
            elif coluna.lower().rsplit(' ')[0] == 'cod.' or coluna.lower().rsplit(' ')[0] == 'codigo' or coluna.lower().rsplit(' ')[0] == 'código' or coluna.lower().rsplit(' ')[0] == 'code' or coluna.lower().rsplit(' ')[0] == 'cod' or coluna.lower().rsplit(' ')[0] == 'cód' or coluna.lower().rsplit(' ')[0] == 'cód.':
                codigo_column = coluna
            elif coluna.lower().rsplit(' ')[0] == 'foto' or coluna.lower().rsplit(' ')[0] == 'imagem' or coluna.lower().rsplit(' ')[0] == 'imagen' or coluna.lower().rsplit(' ')[0] == 'fotografia' or coluna.lower().rsplit(' ')[0] == 'fotografía' or coluna.lower().rsplit(' ')[0] == 'image' or coluna.lower().rsplit(' ')[0] == 'picture' or coluna.lower().rsplit(' ')[0] == 'photo':
                foto_column = coluna
            else:
                other_columns.append(coluna)
        for coluna in [endereco_column, latitude_column, longitude_column, codigo_column, foto_column]:
            if coluna == None:
                if lang == 'es' or lang == 'es-ar':
                    message = f'No se generó el libro {capa["nome"]}. Motivo: Falta la columna obligatoria.'
                elif lang == 'en':
                    message = f'The book {capa["nome"]} was not generated. Reason: Missing required column.'
                else:
                    message = f'O book {capa["nome"]} não foi gerado. Motivo: Ausência de coluna obrigatória.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
                print(f'Erro ao gerar o book {capa["nome"]}. Motivo: ausência de coluna obrigatória.')
                return False

        # Gerando PDF
        pdf_buffer = BytesIO()
        pdf = canvas.Canvas(pdf_buffer, (400*mm, 220*mm))
        # Gerando PPTX
        apresentacao = Presentation()
        apresentacao.slide_height = Mm(220)
        apresentacao.slide_width = Mm(400)

        # Capa
        
        if not capa['nome']:
            message = f'O book não foi gerado. Motivo: Você deve fornecer um nome para o book.'
            send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
            print(f'Erro ao gerar o book. Motivo: Nome para o book não foi fornecido.')
            return False
        else:
            # PDF
            pdf.drawImage('app/static/media/pdf_provider_images/capa_template.jpg', 0, 0, 400*mm, 220*mm)
            pdf.setFont('Helvetica', 10*mm)
            
            # Inicializando slide da capa
            slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
            capa_template = slide.shapes.add_picture('app/static/media/pdf_provider_images/capa_template.jpg', Mm(0), Mm(0), height=Mm(220), width=Mm(400))
            
            if not capa['cliente'] and not capa['pessoa']:
                pdf.drawString(105*mm,25*mm, capa['nome'])

                capa_nome = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                capa_nome_text_frame = capa_nome.text_frame
                capa_nome_text_frame.clear()
                capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                capa_nome_text.text = capa['nome']
                capa_nome_text.font.name = 'Helvetica'
                capa_nome_text.font.size = Mm(10)
                
            else:
                if capa['cliente'] and capa['pessoa']:
                    texto = f'A\C {capa["pessoa"]} - {capa["cliente"]}'
                    if len(texto) > 48:
                        pdf.drawString(105*mm,55*mm, capa['nome'])
                        
                        capa_nome = slide.shapes.add_textbox(Mm(102), Mm(155), Mm(200), Mm(10))
                        capa_nome_text_frame = capa_nome.text_frame
                        capa_nome_text_frame.clear()
                        capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                        capa_nome_text.text = capa['nome']
                        capa_nome_text.font.name = 'Helvetica'
                        capa_nome_text.font.size = Mm(10)

                        pdf.drawString(105*mm,40*mm, texto[0:48])

                        texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                        texto_pptx1_text_frame = texto_pptx1.text_frame
                        texto_pptx1_text_frame.clear()
                        texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                        texto_pptx1_text.text = texto[0:48]
                        texto_pptx1_text.font.name = 'Helvetica'
                        texto_pptx1_text.font.size = Mm(10)

                        pdf.drawString(105*mm,25*mm, texto[48:len(texto)])

                        texto_pptx2 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                        texto_pptx2_text_frame = texto_pptx2.text_frame
                        texto_pptx2_text_frame.clear()
                        texto_pptx2_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        texto_pptx2_text = texto_pptx2_text_frame.paragraphs[0].add_run()
                        texto_pptx2_text.text = texto[48:len(texto)]
                        texto_pptx2_text.font.name = 'Helvetica'
                        texto_pptx2_text.font.size = Mm(10)

                    else:
                        pdf.drawString(105*mm,40*mm, capa['nome'])
                        
                        capa_nome = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                        capa_nome_text_frame = capa_nome.text_frame
                        capa_nome_text_frame.clear()
                        capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                        capa_nome_text.text = capa['nome']
                        capa_nome_text.font.name = 'Helvetica'
                        capa_nome_text.font.size = Mm(10)

                        pdf.drawString(105*mm,25*mm, texto)

                        texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                        texto_pptx1_text_frame = texto_pptx1.text_frame
                        texto_pptx1_text_frame.clear()
                        texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                        texto_pptx1_text.text = texto
                        texto_pptx1_text.font.name = 'Helvetica'
                        texto_pptx1_text.font.size = Mm(10)

                else:
                    pdf.drawString(105*mm,40*mm, capa['nome'])

                    capa_nome = slide.shapes.add_textbox(Mm(102), Mm(170), Mm(200), Mm(10))
                    capa_nome_text_frame = capa_nome.text_frame
                    capa_nome_text_frame.clear()
                    capa_nome_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    capa_nome_text = capa_nome_text_frame.paragraphs[0].add_run()
                    capa_nome_text.text = capa['nome']
                    capa_nome_text.font.name = 'Helvetica'
                    capa_nome_text.font.size = Mm(10)
                    if capa['cliente'] and not capa['pessoa']:
                        texto = f'{capa["cliente"]}'
                        pdf.drawString(105*mm,25*mm, texto)

                        texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                        texto_pptx1_text_frame = texto_pptx1.text_frame
                        texto_pptx1_text_frame.clear()
                        texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                        texto_pptx1_text.text = texto
                        texto_pptx1_text.font.name = 'Helvetica'
                        texto_pptx1_text.font.size = Mm(10)

                    elif capa['pessoa'] and not capa['cliente']:
                        texto = f'A\C {capa["pessoa"]}'
                        pdf.drawString(105*mm,25*mm, texto)

                        texto_pptx1 = slide.shapes.add_textbox(Mm(102), Mm(185), Mm(200), Mm(10))
                        texto_pptx1_text_frame = texto_pptx1.text_frame
                        texto_pptx1_text_frame.clear()
                        texto_pptx1_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        texto_pptx1_text = texto_pptx1_text_frame.paragraphs[0].add_run()
                        texto_pptx1_text.text = texto
                        texto_pptx1_text.font.name = 'Helvetica'
                        texto_pptx1_text.font.size = Mm(10)

        pdf.showPage()

        for i, linha in enumerate(linhas):
            print('gerando linha')
            foto_link = str(linha[foto_column])
            print(foto_link)
            valid_image = True
            if str(foto_link) == 'nan':
                valid_image = False
            elif not 'http' in str(foto_link):
                valid_image = False
            elif foto_link[len(foto_link) - 4:len(foto_link)] not in ['.png', '.jpg', 'jpeg']:
                valid_image = False
            else:
                with open(f'app/static/media/pdf_provider_images/temp_image{i}.png', 'wb') as nova_imagem:
                    imagem = requests.get(foto_link, stream=True, headers={'User-Agent': 'Chrome/108.0.5359.124'})
                    if not imagem.ok:
                        valid_image = False
                    else:
                        for dado in imagem.iter_content():
                            nova_imagem.write(dado)
            if valid_image == True:
                imagem_pil = Image.open(f'app/static/media/pdf_provider_images/temp_image{i}.png')
                size = imagem_pil.size
                new_image = None
                if size[0] > 1280:
                    new_image = imagem_pil.resize((1280, 768))
                    new_image = new_image.convert('RGB')
                else:
                    new_image = imagem_pil.convert('RGB')
                new_image.save(f'app/static/media/pdf_provider_images/temp_image{i}.jpg', 'JPEG')


            coordenadas = f'{linha[latitude_column]}  {linha[longitude_column]}'

            # Inicializando PPT
            slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
            
            # Template de fundo PDF
            pdf.drawImage('app/static/media/pdf_provider_images/template_pdf.jpg', 0, 0, 400*mm, 220*mm)

            # Template de fundo PPTX
            template_fundo = slide.shapes.add_picture('app/static/media/pdf_provider_images/template_pdf.jpg', Mm(0), Mm(0), height=Mm(220), width=Mm(400))

            # Endereço PDF
            pdf.setFont('Helvetica-Bold', 7*mm)
            pdf.setFillColor(colors.white)
            pdf.drawCentredString(190*mm,207*mm, str(linha[endereco_column]))

            # Endereço PPTX
            endereco = slide.shapes.add_textbox(Mm(0), Mm(5), Mm(420), Mm(5))
            endereco_text_frame = endereco.text_frame
            endereco_text_frame.clear()
            endereco_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            endereco_text = endereco_text_frame.paragraphs[0].add_run()
            endereco_text.text = str(linha[endereco_column])
            endereco_text.font.name = 'Helvetica'
            endereco_text.font.bold = True
            endereco_text.font.size = Mm(7)
            endereco_text.font.color.rgb = RGBColor(255,255,255)

            if valid_image == True:
                # Imagem PDF
                pdf.drawImage(f'app/static/media/pdf_provider_images/temp_image{i}.jpg', 3*mm, 18.19*mm, 260*mm, 180*mm)

                #Imagem PPTX
                imagem = slide.shapes.add_picture(f'app/static/media/pdf_provider_images/temp_image{i}.jpg', Mm(3), Mm(21.81), height=Mm(180), width=Mm(262))
            else:
                # Imagem invalida PDF
                pdf.setFont('Helvetica-Bold', 5*mm)
                pdf.setFillColor(colors.black)
                pdf.drawString(48*mm, 80*mm, 'Link inválido ou bloqueado pelo provedor.')

                # Imagem Inválida PPTX
                invalid_img_pptx = slide.shapes.add_textbox(Mm(47), Mm(125), Mm(30), Mm(5))
                invalid_img_text_frame = invalid_img_pptx.text_frame
                invalid_img_text_frame.clear()
                invalid_img_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                invalid_img_text = invalid_img_text_frame.paragraphs[0].add_run()
                invalid_img_text.text = 'Link inválido ou bloqueado pelo provedor.'
                invalid_img_text.font.name = 'Helvetica'
                invalid_img_text.font.bold = True
                invalid_img_text.font.size = Mm(5)

            # Coordenadas PDF
            pdf.setFont('Helvetica-Bold', 5.5*mm)
            pdf.setFillColor(colors.black)
            pdf.drawString(4*mm, 10*mm, 'Coordenadas:')
            
            pdf.setFont('Helvetica', 5.5*mm)
            pdf.drawString(41.5*mm, 10*mm, str(coordenadas))

            # Coordenadas PPTX
            coordenadas_pptx = slide.shapes.add_textbox(Mm(3), Mm(205), Mm(30), Mm(5))
            coordenadas_text_frame = coordenadas_pptx.text_frame
            coordenadas_text_frame.clear()
            coordenadas_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            coordenadas_text = coordenadas_text_frame.paragraphs[0].add_run()
            coordenadas_text.text = 'Coordenadas:'
            coordenadas_text.font.name = 'Helvetica'
            coordenadas_text.font.bold = True
            coordenadas_text.font.size = Mm(5.5)

            coordenadas_content = slide.shapes.add_textbox(Mm(40.5), Mm(205), Mm(50), Mm(6))
            coordenadas_content_text_frame = coordenadas_content.text_frame
            coordenadas_content_text_frame.clear()
            coordenadas_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            coordenadas_content_text = coordenadas_content_text_frame.paragraphs[0].add_run()
            coordenadas_content_text.text = str(coordenadas)
            coordenadas_content_text.font.name = 'Helvetica'
            coordenadas_content_text.font.size = Mm(5.5)
            
            # Código PDF
            pdf.setFont('Helvetica-Bold', 5.5*mm)
            pdf.drawString(180*mm, 10*mm, 'Código:')
            
            pdf.setFont('Helvetica', 5.5*mm)
            pdf.drawString(202*mm, 10*mm, str(linha[codigo_column]))

            # Código PPTX
            codigo = slide.shapes.add_textbox(Mm(180), Mm(205), Mm(30), Mm(5))
            codigo_text_frame = codigo.text_frame
            codigo_text_frame.clear()
            codigo_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            codigo_text = codigo_text_frame.paragraphs[0].add_run()
            codigo_text.text = 'Codigo:'
            codigo_text.font.name = 'Helvetica'
            codigo_text.font.bold = True
            codigo_text.font.size = Mm(5.5)

            codigo_content = slide.shapes.add_textbox(Mm(202), Mm(205), Mm(50), Mm(6))
            codigo_content_text_frame = codigo_content.text_frame
            codigo_content_text_frame.clear()
            codigo_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            codigo_content_text = codigo_content_text_frame.paragraphs[0].add_run()
            codigo_content_text.text = str(linha[codigo_column])
            codigo_content_text.font.name = 'Helvetica'
            codigo_content_text.font.size = Mm(5.5)

            eixo_y_pdf = 194
            eixo_y_pptx = 20

            for coluna in other_columns:
                titulo = str(coluna)
                if titulo == 'nan':
                    titulo = ''
                
                conteudo = ''
                if coluna in linha and str(linha[coluna]) != 'nan':
                    conteudo = str(linha[coluna])
                
                # Se não tiver conteúdo, pula a linha
                if not conteudo.strip() or conteudo.strip() == '':
                    print(titulo)
                    continue
                # Outras Colunas PDF
                pdf.setFont('Helvetica-Bold', 5.5*mm)
                pdf.drawString(267*mm, eixo_y_pdf*mm, titulo)
                
                pdf.setFont('Helvetica', 5.5*mm)
                pdf.drawString(267*mm, (eixo_y_pdf - 7.5) * mm, conteudo)
                eixo_y_pdf -= 17

                # Outras colunas PPTX
                outros = slide.shapes.add_textbox(Mm(266), Mm(eixo_y_pptx), Mm(30), Mm(6))
                outros_text_frame = outros.text_frame
                outros_text_frame.clear()
                outros_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                outros_text_frame.paragraphs[0].word_wrap = True
                outros_text = outros_text_frame.paragraphs[0].add_run()
                outros_text.text = titulo
                outros_text.font.name = 'Helvetica'
                outros_text.font.bold = True
                outros_text.font.size = Mm(5.5)

                outros_content = slide.shapes.add_textbox(Mm(266), Mm(eixo_y_pptx + 7.5), Mm(30), Mm(6))
                outros_content_text_frame = outros_content.text_frame
                outros_content_text_frame.clear()
                outros_content_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                outros_content_text_frame.paragraphs[0].word_wrap = True
                outros_content_text = outros_content_text_frame.paragraphs[0].add_run()
                outros_content_text.text = conteudo
                outros_content_text.font.name = 'Helvetica'
                outros_content_text.font.size = Mm(5.5)
                eixo_y_pptx += 17
            
            # Página PDF
            pdf.setFont('Helvetica-Bold', 5.5*mm)
            pdf.setFillColor(colors.white)
            pdf.drawCentredString(387*mm, 6.5*mm, str(pdf.getPageNumber()))

            # Página PPTX
            pagina = slide.shapes.add_textbox(Mm(380), Mm(203.5), Mm(10), Mm(10))
            pagina_text_frame = pagina.text_frame
            pagina_text_frame.clear()
            pagina_text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            pagina_text = pagina_text_frame.paragraphs[0].add_run()
            pagina_text.text = str(pdf.getPageNumber())
            pagina_text.font.name = 'Helvetica'
            pagina_text.font.bold = True
            pagina_text.font.size = Mm(5.5)
            pagina_text.font.color.rgb = RGBColor(255,255,255)
            
            pdf.showPage()
            
        pdf.save()

        #Gravando pdf
        with pdf_buffer as pdf_file:
            pdf_upload = upload_file_to_s3(pdf_file.getvalue(), f'pdf/{image_id}.pdf', 'application/pdf')
            if pdf_upload == False:
                message = 'Falha ao salvar o arquivo PDF. Apague o book e tente novamente.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
        pdf_buffer.close()
        
        # Gravando pptx
        pptx_buffer = BytesIO()
        apresentacao.save(pptx_buffer)
        with pptx_buffer as pptx_file:
            pptx_upload = upload_file_to_s3(pptx_file.getvalue(), f'pptx/{image_id}.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            if pptx_upload == False:
                message = 'Falha ao salvar o arquivo PPTX. Apague o book e tente novamente.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
        pptx_buffer.close()

        if is_worker == True:
            session.add(Worksheet_Content(
                user_id,
                capa['nome'].title(),
                capa['cliente'],
                capa['pessoa'],
                dumps(content),
                date.today(),
                image_id
            ))
            session.commit()
        message = f'Book {capa["nome"].title()} cadastrado e seus arquivos PDF e PPTX gerados com sucesso.'
        send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
        print('terminando de gerar pdf')
            
    except UnidentifiedImageError as error:
        print(str(error))
        message = f'Falha ao gerar o Book {capa["nome"]}. Motivo: Link de imagem inválido. Você deve fornecer um link de imagem .jpg, .jpeg ou .png. Links de telas de Google Maps ou Street View não são aceitas.'
        send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
    except Exception as error:
        print(str(error))
        message = f'Falha ao gerar o Boook {capa["nome"]}. Motivo: Erro no servidor.'
        send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})

def concatLinhas(linhasList):
    text = ''
    for linha in linhasList:
        text = f'{text} {str(linha)},'
    return text

def points_register(arquivo, nome_arquivo, lang, user_id, pattern_columns, is_worker, db_session=None):
    messages = []
    print('até aqui')
    try:
        tabela = pd.ExcelFile(arquivo)
        planilhas_count = len(tabela.sheet_names)
        if planilhas_count != 1:
            if not is_worker:
                if lang == 'es' or lang == 'es-ar':
                    messages.append('Su hoja de trabajo tiene dos pestañas y solo se procesó la primera.')
                elif lang == 'en':
                    messages.append('Your worksheet has two tabs, and only the first one was processed.')
                else:
                    messages.append('Sua planilha possui duas abas, e somente a primeira foi processada.')
            else:
                if lang == 'es' or lang == 'es-ar':
                    message = 'Su hoja de trabajo tiene dos pestañas y solo se procesó la primera.'
                elif lang == 'en':
                    message = 'Your worksheet has two tabs, and only the first one was processed.'
                else:
                    message = 'Sua planilha possui duas abas, e somente a primeira foi processada.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})

        planilha = pd.read_excel(arquivo, sheet_name=tabela.sheet_names[0]).to_dict('records')
        colunas = pd.read_excel(arquivo, sheet_name=tabela.sheet_names[0])
        code_col, address_col, latitude_col, longitude_col, image_col = None, None, None, None, None
        reference_col, district_col, county_col, zone_col, state_col, country_col, format_col, measure_col = None, None, None, None, None, None, None, None
        impacto_col, valor_tab_comm, valor_negociado_comm, producao_col, observacoes_col, outros_comm_col = None, None, None, None, None, None
        empresa_col, measure_int_col, valor_negociado_int, custo_liq_col, observacoes_int_col, outros_int_col = None, None, None, None, None, None
        for i, coluna in enumerate(colunas):
            if is_in_the_column(coluna.lower(), ['cod', 'cod.', 'cód', 'cód.', 'código', 'codigo', 'code']):
                code_col = coluna
            elif is_in_the_column(coluna.lower(), ['endereço', 'endereco', 'direccion', 'dirección', 'ubicacion', 'ubicación', 'address']):
                address_col = coluna
            elif is_in_the_column(coluna.lower(), ['latitude', 'latitud', 'latitúd']):
                latitude_col = coluna
            elif is_in_the_column(coluna.lower(), ['longitude', 'longitud', 'longitúd']):
                longitude_col = coluna
            elif is_in_the_column(coluna.lower(), ['foto', 'imagem', 'imagen', 'fotografia', 'fotografía', 'image', 'photo', 'picture']):
                image_col = coluna
            elif coluna.lower() in ['referencia', 'referencia ', 'referência', 'referência ', 'reference', 'reference ']:
                reference_col = coluna
            elif coluna.lower() in ['bairro', 'bairro ', 'distrito', 'distrito ', 'barrio', 'barrio ', 'district', 'district ']:
                district_col = coluna
            elif coluna.lower() in ['cidade', 'cidade ', 'ciudad', 'ciudad ', 'city', 'city ', 'county', 'county ']:
                county_col = coluna
            elif coluna.lower() in ['zona', 'zona ', 'zone', 'zone ', 'regiao', 'regiao ', 'região', 'região ', 'región', 'región ', 'region', 'region ']:
                zone_col = coluna
            elif coluna.lower() in ['estado', 'estado ', 'provincia', 'provincia ', 'state', 'state ']:
                state_col = coluna
            elif coluna.lower() in ['país', 'país ', 'pais', 'pais ', 'country', 'country ']:
                country_col = coluna
            elif coluna.lower() in ['formato', 'formato ', 'format', 'format ']:
                format_col = coluna
            elif coluna.lower() in ['medida', 'medida ', 'tamanho', 'tamanho ', 'tamaño', 'tamaño ', 'tamano', 'tamano ', 'measure', 'measure ', 'size', 'size ']:
                measure_col = coluna
            elif coluna.lower() in ['impacto', 'impacto ', 'fluxo passantes', 'fluxo passantes ', 'impact', 'impact ']:
                impacto_col = coluna
            elif coluna.lower() in ['valor tabela', 'valor tabela ', 'valor de la tabla', 'valor de la tabla ', 'average media value', 'average media value ']:
                valor_tab_comm = coluna
            elif coluna.lower() in ['valor negociado', 'valor negociado ', 'our media value', 'our media value ']:
                valor_negociado_comm = coluna
            elif coluna.lower() in ['produção', 'produção ', 'producao', 'producao ', 'produccion', 'produccion ', 'producción', 'producción ', 'production', 'production ']:
                producao_col = coluna
            elif coluna.lower() in ['observações', 'observações ', 'observacoes', 'observacoes ', 'comentarios', 'comentarios ', 'comments', 'comments ']:
                observacoes_col = coluna
            elif coluna.lower() in ['empresa', 'empresa ', 'company', 'company ']:
                empresa_col = coluna
            elif coluna.lower() in ['medida interna', 'medida interna ', 'internal measurement', 'internal measurement ']:
                measure_int_col = coluna
            elif coluna.lower() in ['valor negociado interno', 'valor negociado interno ', 'internal negotiated value', 'internal negotiated value ']:
                valor_negociado_int = coluna
            elif coluna.lower() in ['custo líquido', 'custo líquido ', 'custo liquido', 'custo liquido ', 'preço líquido', 'preço líquido ', 'preco liquido', 'preco liquido ', 'preço liquido', 'preço liquido ', 'preco líquido', 'preco líquido ', 'costo neto', 'costo neto ', 'net cost', 'net cost ']:
                custo_liq_col = coluna
            elif coluna.lower() in ['comentários internos', 'comentários internos ', 'comentarios internos', 'comentarios internos ', 'internal comments', 'internal comments ']:
                observacoes_int_col = coluna
            else:
                continue
        if not code_col or not address_col or not latitude_col or not longitude_col or not image_col:
            if not is_worker:
                if lang == 'es' or lang == 'es-ar':
                    messages.append('No se reconoció alguna columna obligatoria. Las columnas obligatorias son: Código, Dirección, Latitud, Longitud y Foto.')
                    return False, messages
                elif lang == 'en':
                    messages.append('Some mandatory column was not recognized. The mandatory columns are: Code, Address, Latitude, Longitude and Photo.')
                    return False, messages
                else:
                    messages.append('Alguma coluna obrigatória não foi reconhecida. As colunas obrigatórias são: Código, Endereço, Latitude, Longitude e Foto.')
                    return False, messages
            else:
                if lang == 'es' or lang == 'es-ar':
                    message = 'No se reconoció alguna columna obligatoria. Las columnas obligatorias son: Código, Dirección, Latitud, Longitud y Foto.'
                elif lang == 'en':
                    message = 'Some mandatory column was not recognized. The mandatory columns are: Code, Address, Latitude, Longitude and Photo.'
                else:
                    message = 'Alguma coluna obrigatória não foi reconhecida. As colunas obrigatórias são: Código, Endereço, Latitude, Longitude e Foto.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})

        pais = pattern_columns['pais'] if pattern_columns else None
        zona = pattern_columns['zona'] if pattern_columns else None
        cidade = pattern_columns['cidade'] if pattern_columns else None
        estado = pattern_columns['estado'] if pattern_columns else None
        empresa = pattern_columns['empresa'] if pattern_columns else None
        custo = pattern_columns['custo'] if pattern_columns else None
        valorTabela = pattern_columns['valorTabela'] if pattern_columns else None
        valorNegociado = pattern_columns['valorNegociado'] if pattern_columns else None
        formato = pattern_columns['formato'] if pattern_columns else None

        empresas = []
        linhasNaoProcessadas = []
        for i, linha in enumerate(planilha):
            # Buscando lat e lng com api Google Maps caso não haja latitude e longitude cadastrada na linha
            if not str(linha[latitude_col]) or str(linha[latitude_col]) == 'nan' or not str(linha[longitude_col]) or str(linha[longitude_col]) == 'nan':
                temp_endereco = str(linha[address_col])
                temp_bairro = linha[district_col] if district_col and str(linha[district_col]) != 'nan' else ''
                temp_zona = zona if zona else linha[zone_col] if zone_col and str(linha[zone_col]) != 'nan' else ''
                temp_cidade = cidade if cidade else linha[county_col] if county_col and str(linha[county_col]) != 'nan' else ''
                endereco_completo = f'{temp_endereco}, {temp_bairro}, {temp_zona}, {temp_cidade}'
                geocode_result = gmaps.geocode(endereco_completo)
                if len(geocode_result) > 0:
                    linha[latitude_col] = geocode_result[0]['geometry']['location']['lat']
                    linha[longitude_col] = geocode_result[0]['geometry']['location']['lng']
                else:
                    linhasNaoProcessadas.append(i + 2)
                    continue

            if str(linha[code_col]) == 'nan' or str(linha[address_col]) == 'nan' or str(linha[latitude_col]) == 'nan' or str(linha[longitude_col]) == 'nan' or str(linha[image_col]) == 'nan':
                if not is_worker:
                    if lang == 'es' or lang == 'es-ar':
                        messages.append('A la hoja de trabajo le faltan los datos requeridos. Corrija y vuelva a intentarlo.')
                    elif lang == 'en':
                        messages.append('The worksheet is missing required data. Please correct and try again.')
                    else:
                        messages.append('A planilha está faltando dados obrigatórios. Corrija e tente novamente.')
                    return False, messages
                else:
                    if lang == 'es' or lang == 'es-ar':
                        message = f'Faltan datos obligatorios en alguna fila de la hoja de trabajo {nome_arquivo}. Corrija y vuelva a intentarlo.'
                    elif lang == 'en':
                        message = f'Some row of worksheet {nome_arquivo} is missing mandatory data. Please correct and try again.'
                    else:
                        message = f'Alguma linha da planilha {nome_arquivo} está faltando dados obrigatórios. Corrija e tente novamente.'
                    send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
                    return
            if not is_float(linha[latitude_col]) or not is_float(linha[longitude_col]):
                linhasNaoProcessadas.append(i + 2)
                continue
            
            linhaPais = linha[country_col] if country_col and str(linha[country_col]) != 'nan' else None
            linhaZona = linha[zone_col] if zone_col and str(linha[zone_col]) != 'nan' else None
            linhaCidade = linha[county_col] if county_col and str(linha[county_col]) != 'nan' else None
            linhaEstado = linha[state_col] if state_col and str(linha[state_col]) != 'nan' else None
            linhaEmpresa = linha[empresa_col] if empresa_col and str(linha[empresa_col]) != 'nan' else None
            linhaCusto = linha[custo_liq_col] if custo_liq_col and str(linha[custo_liq_col]) != 'nan' else None
            linhaValorTabela = linha[valor_tab_comm] if valor_tab_comm and str(linha[valor_tab_comm]) != 'nan' else None
            linhaValorNegociado = linha[valor_negociado_comm] if valor_negociado_comm and str(linha[valor_negociado_comm]) != 'nan' else None
            linhaFormato = linha[format_col] if format_col and str(linha[format_col]) != 'nan' else None

            empresa_for_register = empresa if empresa else linhaEmpresa
            if empresa_for_register and not empresa_for_register in empresas:
                empresas.append(empresa_for_register)
            if not is_worker:
                new_spot = Flask_Spot(
                    linha[code_col],
                    linha[address_col],
                    float(linha[latitude_col]),
                    float(linha[longitude_col]),
                    linha[image_col],
                    date.today(),
                    linha[reference_col] if reference_col and str(linha[reference_col]) != 'nan' else None,
                    linha[district_col] if district_col and str(linha[district_col]) != 'nan' else None,
                    cidade if cidade else linhaCidade,
                    zona if zona else linhaZona,
                    estado if estado else linhaEstado,
                    pais if pais else linhaPais,
                    formato if formato else linhaFormato,
                    linha[measure_col] if measure_col and str(linha[measure_col]) != 'nan' else None
                )
                new_spot_com = Flask_Spot_Comm(
                    None,
                    linha[impacto_col] if impacto_col and str(linha[impacto_col]) != 'nan' else None,
                    valorTabela if valorTabela else linhaValorTabela,
                    valorNegociado if valorNegociado else linhaValorNegociado,
                    linha[producao_col] if producao_col and str(linha[producao_col]) != 'nan' else None,
                    linha[observacoes_col] if observacoes_col and str(linha[observacoes_col]) != 'nan' else None,
                    outros_comm_col
                )
                new_spot_int = Flask_Spot_Int(
                    None,
                    empresa if empresa else linhaEmpresa,
                    linha[measure_int_col] if measure_int_col and str(linha[measure_int_col]) != 'nan' else None,
                    linha[valor_negociado_int] if valor_negociado_int and str(linha[valor_negociado_int]) != 'nan' else None,
                    custo if custo else linhaCusto,
                    linha[observacoes_int_col] if observacoes_int_col and str(linha[observacoes_int_col]) != 'nan' else None,
                    outros_int_col
                )
                if db_session == None:
                    db.session.add(new_spot)
                    db.session.flush()
                    db.session.refresh(new_spot)

                    new_spot_com.spot_id = new_spot.id
                    db.session.add(new_spot_com)
                    
                    new_spot_int.spot_id = new_spot.id
                    db.session.add(new_spot_int)
                else:
                    db_session.add(new_spot)
                    db_session.flush()
                    db_session.refresh(new_spot)

                    new_spot_com.spot_id = new_spot.id
                    db_session.add(new_spot_com)
                    
                    new_spot_int.spot_id = new_spot.id
                    db_session.add(new_spot_int)
            else:
                new_spot = Spot(
                    linha[code_col],
                    linha[address_col],
                    float(linha[latitude_col]),
                    float(linha[longitude_col]),
                    linha[image_col],
                    date.today(),
                    linha[reference_col] if reference_col and str(linha[reference_col]) != 'nan' else None,
                    linha[district_col] if district_col and str(linha[district_col]) != 'nan' else None,
                    cidade if cidade else linhaCidade,
                    zona if zona else linhaZona,
                    estado if estado else linhaEstado,
                    pais if pais else linhaPais,
                    formato if formato else linhaFormato,
                    linha[measure_col] if measure_col and str(linha[measure_col]) != 'nan' else None
                )
                new_spot_com = Spot_Commercial_Info(
                    None,
                    linha[impacto_col] if impacto_col and str(linha[impacto_col]) != 'nan' else None,
                    valorTabela if valorTabela else linhaValorTabela,
                    valorNegociado if valorNegociado else linhaValorNegociado,
                    linha[producao_col] if producao_col and str(linha[producao_col]) != 'nan' else None,
                    linha[observacoes_col] if observacoes_col and str(linha[observacoes_col]) != 'nan' else None,
                    outros_comm_col
                )
                new_spot_int = Spot_Private_Info(
                    None,
                    empresa if empresa else linhaEmpresa,
                    linha[measure_int_col] if measure_int_col and str(linha[measure_int_col]) != 'nan' else None,
                    linha[valor_negociado_int] if valor_negociado_int and str(linha[valor_negociado_int]) != 'nan' else None,
                    custo if custo else linhaCusto,
                    linha[observacoes_int_col] if observacoes_int_col and str(linha[observacoes_int_col]) != 'nan' else None,
                    outros_int_col
                )
                session.add(new_spot)
                session.flush()
                session.refresh(new_spot)

                new_spot_com.spot_id = new_spot.id
                session.add(new_spot_com)
                
                new_spot_int.spot_id = new_spot.id
                session.add(new_spot_int)
        
        if not is_worker:
            fornecedores = get_fornecedores_list()
            fornecedores_novos = []
            for empresa in empresas:
                if not empresa in fornecedores:
                    fornecedores_novos.append(empresa)
                
            if db_session == None:
                if len(fornecedores_novos) > 0:
                    for fornecedor in fornecedores_novos:
                        new = Flask_Person(fornecedor, None, None, None, None, None, None, None, 1)
                        db.session.add(new)
                db.session.commit()
            else:
                if len(fornecedores_novos) > 0:
                    for fornecedor in fornecedores_novos:
                        new = Flask_Person(fornecedor, None, None, None, None, None, None, None, 1)
                        db_session.add(new)

            if len(linhasNaoProcessadas) > 0:
                if lang == 'es' or lang == 'es-ar':
                    messages.append(f'Las lineas{concatLinhas(linhasNaoProcessadas)} no se procesaron, ya que no se encontraron datos de latitud y longitud.')
                elif lang == 'en':
                    messages.append(f'Lines{concatLinhas(linhasNaoProcessadas)} were not processed, as no latitude and longitude data were found.')
                else:
                    messages.append(f'As linhas{concatLinhas(linhasNaoProcessadas)} não foram processadas, pois não foram encontrados dados de Latitude e Longitude.')
            if lang == 'es' or lang == 'es-ar':
                messages.append('Los puntos se han registrado con éxito.')
            elif lang == 'en':
                messages.append('Points have been registered successfully.')
            else:
                messages.append('Os pontos foram cadastrados com sucesso.')
            return True, messages
        else:
            fornecedores = get_fornecedores_list_off_context()
            fornecedores_novos = []
            for empresa in empresas:
                if not empresa in fornecedores:
                    fornecedores_novos.append(empresa)
            if len(fornecedores_novos) > 0:
                for fornecedor in fornecedores_novos:
                    new = Person(fornecedor, None, None, None, None, None, None, None, 1)
                    session.add(new)
            session.commit()
            session.close()
            if lang == 'es' or lang == 'es-ar':
                message = f'La hoja de cálculo {nome_arquivo} se registró correctamente.'
            elif lang == 'en':
                message = f'Spreadsheet {nome_arquivo} was successfully registered.'
            else:
                message = f'A planilha {nome_arquivo} foi registrada com sucesso.'
            if len(linhasNaoProcessadas) > 0:
                if lang == 'es' or lang == 'es-ar':
                    message = message + f' Las lineas{concatLinhas(linhasNaoProcessadas)} no se procesaron, ya que no se encontraron datos validos de latitud y longitud.'
                elif lang == 'en':
                    message = message + f' Lines{concatLinhas(linhasNaoProcessadas)} were not processed, as no latitude and longitude valid data were found.'
                else:
                    message = message + f' As linhas{concatLinhas(linhasNaoProcessadas)} não foram processadas, pois não foram encontrados válidos dados de Latitude e Longitude.'
            send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
            return
        
    except ValueError as error:
        print(f'erro = {str(error)}')
        if 'Excel file format cannot be determined, you must specify an engine manually.' in str(error):
            if not is_worker:
                if lang == 'es' or lang == 'es-ar':
                    messages.append('Formato de archivo inválido.')
                    return False, messages
                elif lang == 'en':
                    messages.append('Invalid file format.')
                    return False, messages
                else:
                    messages.append('Formato de arquivo inválido.')
                    return False, messages
            else:
                if lang == 'es' or lang == 'es-ar':
                    message = f'La hoja de cálculo {nome_arquivo} no se registró. Motivo: formato de archivo no válido.'
                elif lang == 'en':
                    message = f'Spreadsheet {nome_arquivo} was not registered. Reason: Invalid file format.'
                else:
                    message = f'A planilha {nome_arquivo} não foi cadastrada. Motivo: Formato de arquivo inválido.'
                
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
                return
        else:
            if not is_worker:
                return False, ['Erro no servidor. Tente novamente mais tarde.']
            else:
                message = f'A planilha {nome_arquivo} não foi cadastrada. Erro no servidor.'
                send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
                return
    except Exception as error:
        print(f'erro = {str(error)}')
        if not is_worker:
            print(f'erro = {str(error)}')
            return False, ['Erro no servidor. Tente novamente mais tarde.']
        else:
            message = f'A planilha {nome_arquivo} não foi cadastrada. Erro no servidor.'
            send_message = requests.get(f'{os.environ["APP_URL"]}/flash-message-generate?message={message}&user={user_id}', headers={'Secret-Key': os.environ['SECRET_KEY']})
            return

def book_register(arquivo, dados_capa, colunas, lang, user_id):
    inicio = time.time()
    messages = []
    if len(colunas) > 15:
        if lang == 'es' or lang == 'es-ar':
            messages.append('El libro no se generó porque tiene más de 15 columnas.')
            return False, messages
        elif lang == 'en':
            messages.append('The book was not generated because it has more than 15 columns.')
            return False, messages
        else:
            messages.append('O book não foi gerado porque tem mais de 15 colunas.')
            return False, messages
    tabela = pd.ExcelFile(arquivo)
    planilhas = tabela.sheet_names # Pega o nome de todas as abas da tabela
    for i, planilha in enumerate(planilhas):
        if i > 0:
            continue
        
        # Verifica se as colunas Endereço e Foto estão presentes
        endereco_check = False
        foto_check = False
        latitude_check = False
        longitude_check = False
        codigo_check = False
        for i, coluna in enumerate(colunas):
            if 'endereço' in str(coluna).lower() or 'endereco' in str(coluna).lower() or 'direccion' in str(coluna).lower() or 'dirección' in str(coluna).lower() or 'address' in str(coluna).lower():
                endereco_check = True
            elif 'foto' in str(coluna).lower() or 'imagem' in str(coluna).lower() or 'imagen' in str(coluna).lower() or 'fotografia' in str(coluna).lower() or 'fotografía' in str(coluna).lower() or 'image' in str(coluna).lower() or 'picture' in str(coluna).lower() or 'photo' in str(coluna).lower():
                foto_check = True
            elif 'latitude' in str(coluna).lower() or 'latitud' in str(coluna).lower():
                latitude_check = True
            elif 'longitude' in str(coluna).lower() or 'longitud' in str(coluna).lower():
                longitude_check = True
            elif 'codigo' in str(coluna).lower() or 'código' in str(coluna).lower() or 'cod' == str(coluna).lower() or 'cód' == str(coluna).lower() or 'code' in str(coluna).lower() or 'cod.' in str(coluna).lower() or 'cód.' in str(coluna).lower():
                codigo_check = True
            
        if not foto_check or not endereco_check or not latitude_check or not longitude_check or not codigo_check:
            if lang == 'es' or lang == 'es-ar':
                messages.append('Alguna columna requerida no fue reconocida. Las columnas obligatorias son "Código", "Dirección", "Latitud", "Longitud" y "Foto".')
                return False, messages
            elif lang == 'en':
                messages.append('Some required column was not recognized. Required columns are "Code", "Address", "Latitude", "Longitude" and "Photo".')
                return False, messages
            else:
                messages.append('Alguma coluna obrigatória não foi reconhecida. As colunas obrigatórias são "Código", "Endereço", "Latitude", "Longitude" e "Foto".')
                return False, messages
        
        book = pd.read_excel(arquivo, sheet_name=planilha, nrows=500).to_dict('records')
        new_book = []
        for linha in book:
            if str(linha[colunas[0]]) == 'nan' and\
            str(linha[colunas[1]]) == 'nan' and\
            str(linha[colunas[2]]) == 'nan' and\
            str(linha[colunas[3]]) == 'nan' and\
            str(linha[colunas[4]]) == 'nan':
                # Pula linhas em branco
                continue
            new_book.append(linha)
        image_id = image_id_generator()
        capa = {'nome': dados_capa['nome'], 'cliente': dados_capa['cliente'], 'pessoa': dados_capa['pessoa']}
        content = {'colunas': colunas, 'conteudo': new_book}
        
        result = q.enqueue(pdf_generator, capa, content, image_id, lang, user_id, True, job_timeout='30m')
        if lang == 'es' or lang == 'es-ar':
            messages.append(f'El libro {(dados_capa["nome"]).title()} se está generando. Cuando se complete, aparecerá en "Lista de Books" o notificará en pantalla en caso de error.')
        elif lang == 'en':
            messages.append(f'The book {(dados_capa["nome"]).title()} is being generated. When completed, it will appear in "List of Books" or notify on screen in case of error.')
        else:
            messages.append(f'O book {(dados_capa["nome"]).title()} está sendo gerado. Quando estiver concluído, ele aparecerá em "Lista de Books" ou notificará na tela em caso de erro.')
        fim = time.time()
        print(fim - inicio)
        return True, messages
    
def bookGenerateWithPoints(lang, columns, idList, userId, bookName, personName, clientName):
    messages = []
    try:
        if len(columns) > 15:
            if lang == 'es' or lang == 'es-ar':
                messages.append('Elija un máximo de 15 columnas.')
            elif lang == 'en':
                messages.append('Choose a maximum of 15 columns.')
            else:
                messages.append('Escolha no máximo 15 colunas.')
            response = Register_Response_(
                status=False,
                message=messages
            )
            return response.json()
        allColumns = columns.copy()
        codigo_col = None
        latitude_col = None
        longitude_col = None
        endereco_col = None
        image_col = None
        for coluna in columns:
            if is_in_the_column(coluna.lower(), ['cod', 'cod.', 'cód', 'cód.', 'código', 'codigo', 'code']):
                codigo_col = coluna
            elif is_in_the_column(coluna.lower(), ['endereço', 'endereco', 'direccion', 'dirección', 'ubicacion', 'ubicación', 'address']):
                endereco_col = coluna
            elif is_in_the_column(coluna.lower(), ['latitude', 'latitud', 'latitúd']):
                latitude_col = coluna
            elif is_in_the_column(coluna.lower(), ['longitude', 'longitud', 'longitúd']):
                longitude_col = coluna
            elif is_in_the_column(coluna.lower(), ['foto', 'imagem', 'imagen', 'fotografia', 'fotografía', 'image', 'photo', 'picture']):
                image_col = coluna
        for coluna in [codigo_col, latitude_col, longitude_col, endereco_col, image_col]:
            if coluna:
                allColumns.remove(coluna)
            else:
                if lang == 'es' or lang == 'es-ar':
                    messages.append('Erro en servidor.')
                elif lang == 'en':
                    messages.append('Server error.')
                else:
                    messages.append('Erro no servidor.')
                response = Register_Response_(
                    status=False,
                    message=messages
                )
                return response.json()
        
        pontos = get_pontos_by_id_list(idList)
        conteudo = []
        for ponto in pontos:
            linha = {}
            linha[codigo_col] = ponto[0].code
            linha[latitude_col] = ponto[0].latitude
            linha[longitude_col] = ponto[0].longitude
            linha[endereco_col] = ponto[0].address
            linha[image_col] = ponto[0].image_link
            for coluna in allColumns:
                if coluna.lower() in ['referencia', 'referencia ', 'referência', 'referência ', 'reference', 'reference ']:
                    linha[coluna] = ponto[0].reference if ponto[0].reference else ''
                elif coluna.lower() in ['bairro', 'bairro ', 'distrito', 'distrito ', 'barrio', 'barrio ', 'district', 'district ']:
                    linha[coluna] = ponto[0].district if ponto[0].district else ''
                elif coluna.lower() in ['cidade', 'cidade ', 'ciudad', 'ciudad ', 'city', 'city ', 'county', 'county ']:
                    linha[coluna] = ponto[0].city if ponto[0].city else ''
                elif coluna.lower() in ['zona', 'zona ', 'zone', 'zone ', 'regiao', 'regiao ', 'região', 'região ', 'región', 'región ', 'region', 'region ']:
                    linha[coluna] = ponto[0].zone if ponto[0].zone else ''
                elif coluna.lower() in ['estado', 'estado ', 'provincia', 'provincia ', 'state', 'state ']:
                    linha[coluna] = ponto[0].state if ponto[0].state else ''
                elif coluna.lower() in ['país', 'país ', 'pais', 'pais ', 'country', 'country ']:
                    linha[coluna] = ponto[0].country if ponto[0].country else ''
                elif coluna.lower() in ['formato', 'formato ', 'format', 'format ']:
                    linha[coluna] = ponto[0].format if ponto[0].format else ''
                elif coluna.lower() in ['medida', 'medida ', 'tamanho', 'tamanho ', 'tamaño', 'tamaño ', 'tamano', 'tamano ', 'measure', 'measure ', 'size', 'size ']:
                    linha[coluna] = ponto[0].measure if ponto[0].measure else ''
                elif coluna.lower() in ['impacto', 'impacto ', 'fluxo passantes', 'fluxo passantes ', 'impact', 'impact ']:
                    linha[coluna] = ponto[1].impacto if ponto[1].impacto else ''
                elif coluna.lower() in ['valor tabela', 'valor tabela ', 'valor de la tabla', 'valor de la tabla ', 'average media value', 'average media value ']:
                    linha[coluna] = ponto[1].valor_tabela_comm if ponto[1].valor_tabela_comm else ''
                elif coluna.lower() in ['valor negociado', 'valor negociado ', 'our media value', 'our media value ']:
                    linha[coluna] = ponto[1].valor_negociado_comm if ponto[1].valor_negociado_comm else ''
                elif coluna.lower() in ['produção', 'produção ', 'producao', 'producao ', 'produccion', 'produccion ', 'producción', 'producción ', 'production', 'production ']:
                    linha[coluna] = ponto[1].producao if ponto[1].producao else ''
                elif coluna.lower() in ['observações', 'observações ', 'observacoes', 'observacoes ', 'comentarios', 'comentarios ', 'comments', 'comments ']:
                    linha[coluna] = ponto[1].observacoes if ponto[1].observacoes else ''
                elif coluna.lower() in ['empresa', 'empresa ', 'company', 'company ']:
                    linha[coluna] = ponto[2].empresa if ponto[2].empresa else ''
                elif coluna.lower() in ['medida interna', 'medida interna ', 'internal measurement', 'internal measurement ']:
                    linha[coluna] = ponto[2].medida_int if ponto[2].medida_int else ''
                elif coluna.lower() in ['valor negociado interno', 'valor negociado interno ', 'internal negotiated value', 'internal negotiated value ']:
                    linha[coluna] = ponto[2].valor_negociado_int if ponto[2].valor_negociado_int else ''
                elif coluna.lower() in ['custo líquido', 'custo líquido ', 'custo liquido', 'custo liquido ', 'preço líquido', 'preço líquido ', 'preco liquido', 'preco liquido ', 'preço liquido', 'preço liquido ', 'preco líquido', 'preco líquido ', 'costo neto', 'costo neto ', 'net cost', 'net cost ']:
                    linha[coluna] = ponto[2].custo_liq if ponto[2].custo_liq else ''
                elif coluna.lower() in ['comentários internos', 'comentários internos ', 'comentarios internos', 'comentarios internos ', 'internal comments', 'internal comments ']:
                    linha[coluna] = ponto[2].observacoes if ponto[2].observacoes else ''
                else:
                    continue
            conteudo.append(linha)
        
        capa = {'nome': bookName, 'cliente': clientName, 'pessoa': personName}
        content = {'colunas': columns, 'conteudo': conteudo}
        image_id = image_id_generator()
        
        # Criando fila rq
        q = Queue(connection=conn)

        result = q.enqueue(pdf_generator, capa, content, image_id, lang, userId, True, job_timeout='30m')
        if lang == 'es' or lang == 'es-ar':
            messages.append(f'El libro {(capa["nome"]).title()} se está generando. Cuando se complete, aparecerá en "Lista de Books" o notificará en pantalla en caso de error.')
        elif lang == 'en':
            messages.append(f'The book {(capa["nome"]).title()} is being generated. When completed, it will appear in "List of Books" or notify on screen in case of error.')
        else:
            messages.append(f'O book {(capa["nome"]).title()} está sendo gerado. Quando estiver concluído, ele aparecerá em "Lista de Books" ou notificará na tela em caso de erro.')
        response = Register_Response_(
            status=True,
            message=messages
        )
        return response.json()
    except Exception as error:
        print(str(error))
        if lang == 'es' or lang == 'es-ar':
            messages.append('Erro en servidor.')
        elif lang == 'en':
            messages.append('Server error.')
        else:
            messages.append('Erro no servidor.')
        response = Register_Response_(
            status=False,
            message=messages
        )
        return response.json()

def editar_grupo(lang, idList, dados):
    messages = []
    pontos = get_pontos_by_id_list(idList)
    for ponto in pontos:
        basic = ponto[0]
        commercial = ponto[1]
        private = ponto[2]

        basic.reference = dados['reference'] if dados['reference'] else basic.reference
        basic.district = dados['district'] if dados['district'] else basic.district
        basic.city = dados['city'] if dados['city'] else basic.city
        basic.zone = dados['zone'] if dados['zone'] else basic.zone
        basic.state = dados['state'] if dados['state'] else basic.state
        basic.country = dados['country'] if dados['country'] else basic.country
        basic.format = dados['format'] if dados['format'] else basic.format
        basic.measure = dados['measure'] if dados['measure'] else basic.measure

        commercial.impacto = dados['impacto'] if dados['impacto'] else commercial.impacto
        commercial.valor_tabela_comm = dados['valor_tabela_comm'] if dados['valor_tabela_comm'] else commercial.valor_tabela_comm
        commercial.valor_negociado_comm = dados['valor_negociado_comm'] if dados['valor_negociado_comm'] else commercial.valor_negociado_comm
        commercial.producao = dados['producao'] if dados['producao'] else commercial.producao
        commercial.observacoes = dados['observacoes_comm'] if dados['observacoes_comm'] else commercial.observacoes
        
        private.empresa = dados['empresa'] if dados['empresa'] else private.empresa
        private.medida_int = dados['medida_int'] if dados['medida_int'] else private.medida_int
        private.valor_negociado_int = dados['valor_negociado_int'] if dados['valor_negociado_int'] else private.valor_negociado_int
        private.custo_liq = dados['custo_liq'] if dados['custo_liq'] else private.custo_liq
        private.observacoes = dados['observacoes_int'] if dados['observacoes_int'] else private.observacoes

        db.session.add(basic)
        db.session.add(commercial)
        db.session.add(private)
    db.session.commit()
    if lang == 'es' or lang == 'es-ar':
        messages.append('Puntos cambiados con exito.')
    elif lang == 'en':
        messages.append('Points changed successfully.')
    else:
        messages.append('Pontos alterados com sucesso.')
    response = Register_Response_(
        status=True,
        message=messages
    )
    return response.json()

def gerar_excel(ids, lang):
    messages = []
    pontos = get_pontos_by_id_list(ids)
    planilha = Workbook()
    planilha.remove(planilha['Sheet'])

    nome_planilha = token_hex(32)
    tabela = planilha.create_sheet('POINTS')
    colunas = None
    if lang == 'es' or lang == 'es-ar':
        colunas = [
            'Código',
            'Dirección',
            'Latidud',
            'Longitud',
            'Foto',
            'Barrio',
            'Referencia',
            'Ciudad',
            'Zona',
            'Provincia',
            'País',
            'Formato',
            'Medida',
            'Impacto',
            'Valor de la tabla',
            'Valor negociado',
            'Producción',
            'Comentarios',
            'Empresa',
            'Valor negociado interno',
            'Costo neto',
            'Medida interna',
            'Comentarios internos'
        ]
    elif lang  == 'en':
        colunas = [
            'Code',
            'Address',
            'Latidude',
            'Longitude',
            'Photo',
            'District',
            'Reference',
            'City',
            'Zone',
            'State',
            'Country',
            'Format',
            'Measure',
            'Impact',
            'Average media value',
            'Our media value',
            'Production',
            'Comments',
            'Company',
            'Internal negotiated value',
            'Net cost',
            'Internal measurement',
            'Internal comments'
        ]
    else:
        colunas = [
            'Código',
            'Endereço',
            'Latidude',
            'Longitude',
            'Foto',
            'Bairro',
            'Referência',
            'Cidade',
            'Zona',
            'Estado',
            'País',
            'Formato',
            'Medida',
            'Impacto',
            'Valor tabela',
            'Valor negociado',
            'Produção',
            'Observações',
            'Empresa',
            'Valor negociado interno',
            'Custo líquido',
            'Medida interna',
            'Comentários internos'
        ]
    tabela.append(colunas)
    for ponto in pontos:
        linha = []
        linha.append(ponto[0].code)
        linha.append(ponto[0].address)
        linha.append(ponto[0].latitude)
        linha.append(ponto[0].longitude)
        linha.append(ponto[0].image_link)
        linha.append(ponto[0].district)
        linha.append(ponto[0].reference)
        linha.append(ponto[0].city)
        linha.append(ponto[0].zone)
        linha.append(ponto[0].state)
        linha.append(ponto[0].country)
        linha.append(ponto[0].format)
        linha.append(ponto[0].measure)
        linha.append(ponto[1].impacto)
        linha.append(ponto[1].valor_tabela_comm)
        linha.append(ponto[1].valor_negociado_comm)
        linha.append(ponto[1].producao)
        linha.append(ponto[1].observacoes)
        linha.append(ponto[2].empresa)
        linha.append(ponto[2].valor_negociado_int)
        linha.append(ponto[2].custo_liq)
        linha.append(ponto[2].medida_int)
        linha.append(ponto[2].observacoes)
        tabela.append(linha)
    
    planilha_buffer = BytesIO()
    planilha.save(planilha_buffer)
    with planilha_buffer as planilha_file:
        xlsx_upload = upload_file_to_s3(planilha_file.getvalue(), f'xlsx/{nome_planilha}.xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        if xlsx_upload == False:
            message = 'Falha ao converter arquivo. Tente novamente'
            messages.append(message)
        else:
            if lang == 'es' or lang == 'es-ar':
                messages.append('Hoja de cálculo generada con éxito.')
            elif lang == 'en':
                messages.append('Spreadsheet generated successfully.')
            else:
                messages.append('Planilha gerada com sucesso.')
    planilha_buffer.close()
    response = Register_Response_(
        status=True,
        message=messages,
        data=nome_planilha
    )
    return response.json()
